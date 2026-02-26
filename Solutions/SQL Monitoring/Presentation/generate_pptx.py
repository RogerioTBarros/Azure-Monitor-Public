"""
Generate SQL Server Monitoring Solution PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Load default template from a known-good path (workaround OneDrive path issues)
import shutil
_template_src = r"C:\Users\rogeriob\OneDrive - Microsoft\repos\.venv\Lib\site-packages\pptx\templates\default.pptx"
_template_tmp = r"C:\temp\default.pptx"
shutil.copy2(_template_src, _template_tmp)
prs = Presentation(_template_tmp)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
LIGHT_BLUE = RGBColor(0xDE, 0xEC, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xD1, 0x34, 0x38)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
    return tf

def add_rounded_rect(slide, left, top, width, height, color, text="", font_size=12, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = "Segoe UI"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)
add_shape_bg(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.08), AZURE_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "SQL Server Monitoring Solution", 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Centralized Monitoring with Azure Automation & Logs Ingestion API", 24, False, RGBColor(0xBB, 0xDE, 0xFB), PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "Azure Monitor  |  Log Analytics  |  Custom Workbook", 18, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             "Microsoft Azure Monitor Assets", 16, False, RGBColor(0x64, 0xB5, 0xF6), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "Agenda", 36, True, WHITE)

items = [
    "1.  Challenge & Business Problem",
    "2.  Solution Overview",
    "3.  Architecture Deep Dive",
    "4.  Data Pipeline: Logs Ingestion API",
    "5.  Workbook Dashboard (4 Tabs)",
    "6.  Security & Authentication",
    "7.  Deployment Options",
    "8.  Demo & Next Steps"
]
add_bullet_list(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5), items, 22, DARK_GRAY)

# ============================================================
# SLIDE 3: The Challenge
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "The Challenge", 36, True, WHITE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(12), Inches(0.6),
             "Monitoring SQL Server instances across hybrid environments is complex:", 20, False, DARK_GRAY)

# Challenge boxes
challenges = [
    ("Fragmented Visibility", "SQL Servers spread across on-prem,\nIaaS VMs, and Arc-enabled servers\nwith no unified view"),
    ("Backup Compliance", "No centralized way to verify backup\nSLA compliance across all databases\nand instances"),
    ("Manual Processes", "Teams rely on manual scripts or\nthird-party tools with complex\nlicensing and overhead"),
    ("Reactive Alerting", "Issues discovered after impact;\nno proactive monitoring of\nuptime and backup status")
]

for i, (title, desc) in enumerate(challenges):
    x = Inches(0.8 + i * 3.1)
    box = add_rounded_rect(slide, x, Inches(2.6), Inches(2.8), Inches(3.5), LIGHT_BLUE)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.6), title, 18, True, DARK_BLUE, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(2.2), desc, 14, False, DARK_GRAY, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Solution Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "Solution Overview", 36, True, WHITE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(12), Inches(0.5),
             "A fully native Azure Monitor solution — no additional agents required on SQL Servers", 20, True, AZURE_BLUE)

benefits = [
    ("Agentless", "No software installed\non SQL Servers.\nOnly the Hybrid Worker\nneeds an extension.", GREEN),
    ("Custom Schema", "Clean 20-column table\nin Log Analytics.\nNo JSON parsing needed\nin KQL queries.", AZURE_BLUE),
    ("Secure by Design", "Managed Identity auth.\nKey Vault for credentials.\nNo passwords in code.", DARK_BLUE),
    ("Pre-built Dashboard", "4-tab Azure Monitor\nWorkbook: Summary,\nInstances, Databases,\nBackups.", ORANGE),
]

for i, (title, desc, color) in enumerate(benefits):
    x = Inches(0.8 + i * 3.1)
    add_rounded_rect(slide, x, Inches(2.5), Inches(2.8), Inches(0.6), color, title, 16, WHITE)
    add_text_box(slide, x + Inches(0.2), Inches(3.3), Inches(2.4), Inches(2.5), desc, 14, False, DARK_GRAY, PP_ALIGN.CENTER)

# Key metrics collected
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.4),
             "Metrics Collected:", 16, True, DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.9), Inches(12), Inches(1),
             "Instance Uptime  •  Database State & Recovery Model  •  Full/Log Backup Status  •  Backup SLA Compliance  •  Connection Errors",
             15, False, MEDIUM_GRAY)

# ============================================================
# SLIDE 5: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "Architecture", 36, True, WHITE)

# On-Prem Zone
add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(3), Inches(5.2), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(0.6), Inches(1.55), Inches(2.8), Inches(0.4), "On-Premises / IaaS", 14, True, RGBColor(0x15, 0x65, 0xC0))

for i, name in enumerate(["SQL Server 1", "SQL Server 2", "SQL Server N"]):
    y = Inches(2.2 + i * 1.4)
    add_rounded_rect(slide, Inches(0.8), y, Inches(2.4), Inches(0.8), RGBColor(0x15, 0x65, 0xC0), name, 13, WHITE)

# Hybrid Worker Zone
add_shape_bg(slide, Inches(4), Inches(1.5), Inches(3), Inches(5.2), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(4.1), Inches(1.55), Inches(2.8), Inches(0.4), "Hybrid Worker (Arc VM)", 14, True, RGBColor(0xE6, 0x51, 0x00))

add_rounded_rect(slide, Inches(4.3), Inches(2.2), Inches(2.4), Inches(0.7), RGBColor(0xE6, 0x51, 0x00), "Arc-enabled Server\n+ Managed Identity", 11, WHITE)
add_rounded_rect(slide, Inches(4.3), Inches(3.2), Inches(2.4), Inches(0.7), RGBColor(0xFF, 0xB7, 0x4D), "Hybrid Worker\nExtension (PS 7.2)", 11, DARK_GRAY)
add_rounded_rect(slide, Inches(4.3), Inches(4.2), Inches(2.4), Inches(1.0), RGBColor(0xFF, 0xF9, 0xC4), "PowerShell Runbook\nGet-SQLServerInfo-\nLogsIngestionApi.ps1", 10, DARK_GRAY)

# Azure Zone
add_shape_bg(slide, Inches(7.5), Inches(1.5), Inches(5.3), Inches(5.2), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(7.6), Inches(1.55), Inches(5.1), Inches(0.4), "Azure Cloud", 14, True, RGBColor(0x2E, 0x7D, 0x32))

add_rounded_rect(slide, Inches(7.8), Inches(2.2), Inches(2.2), Inches(0.7), RGBColor(0x2E, 0x7D, 0x32), "Automation\nAccount", 12, WHITE)
add_rounded_rect(slide, Inches(10.3), Inches(2.2), Inches(2.2), Inches(0.7), RGBColor(0xF9, 0xA8, 0x25), "Key Vault\n(optional)", 12, WHITE)
add_rounded_rect(slide, Inches(7.8), Inches(3.3), Inches(2.2), Inches(0.6), RGBColor(0x43, 0xA0, 0x47), "DCE (Endpoint)", 12, WHITE)
add_rounded_rect(slide, Inches(10.3), Inches(3.3), Inches(2.2), Inches(0.6), RGBColor(0x43, 0xA0, 0x47), "DCR (Rule)", 12, WHITE)
add_rounded_rect(slide, Inches(7.8), Inches(4.4), Inches(2.2), Inches(0.9), RGBColor(0x1B, 0x5E, 0x20), "Log Analytics\nSQLServerMonitoring_CL", 11, WHITE)
add_rounded_rect(slide, Inches(10.3), Inches(4.4), Inches(2.2), Inches(0.9), RGBColor(0x6A, 0x1B, 0x9A), "Azure Monitor\nWorkbook (4 tabs)", 11, WHITE)

# Flow arrows (text-based)
add_text_box(slide, Inches(3.2), Inches(2.8), Inches(1.2), Inches(0.5), "TCP 1433 →", 11, True, RGBColor(0x15, 0x65, 0xC0), PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.5), Inches(3.4), Inches(1.2), Inches(0.5), "HTTPS →", 11, True, RGBColor(0x2E, 0x7D, 0x32), PP_ALIGN.CENTER)

# Data Flow
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.4),
             "Data Flow:  ① Schedule triggers runbook  →  ② Hybrid Worker queries SQL Servers  →  ③ POST to Logs Ingestion API  →  ④ DCR routes to custom table  →  ⑤ Workbook visualizes",
             13, False, MEDIUM_GRAY)

# ============================================================
# SLIDE 6: Data Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "Data Pipeline — Logs Ingestion API", 36, True, WHITE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(12), Inches(0.5),
             "Direct ingestion into Log Analytics using Azure Monitor's native REST API", 20, False, AZURE_BLUE)

# Pipeline steps
steps = [
    ("1. Collect", "Runbook queries SQL\nServers using T-SQL\n(sys.dm_os_sys_info,\nsys.databases,\nmsdb.dbo.backupset)", AZURE_BLUE),
    ("2. Transform", "PowerShell converts\nresults to JSON records\nwith 20 typed columns\n(one record per database)", RGBColor(0x43, 0xA0, 0x47)),
    ("3. Authenticate", "Managed Identity\nobtains OAuth2 token\nfor monitor.azure.com\n(no stored credentials)", RGBColor(0xE6, 0x51, 0x00)),
    ("4. Ingest", "HTTPS POST to DCE\nLogs Ingestion API\nwith JSON payload\n(batched per collection)", RGBColor(0x6A, 0x1B, 0x9A)),
    ("5. Route", "DCR validates schema,\napplies transformKql,\nroutes to destination\ntable in Log Analytics", RGBColor(0x2E, 0x7D, 0x32)),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    add_rounded_rect(slide, x, Inches(2.5), Inches(2.2), Inches(0.6), color, title, 14, WHITE)
    add_text_box(slide, x + Inches(0.1), Inches(3.3), Inches(2.0), Inches(2.5), desc, 13, False, DARK_GRAY, PP_ALIGN.CENTER)

# Table schema
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.4),
             "Custom Table: SQLServerMonitoring_CL (20 columns)", 16, True, DARK_BLUE)
schema_text = ("TimeGenerated  •  CollectorName  •  SqlInstance  •  ServerName  •  SqlVersion  •  InstanceStartTime\n"
               "InstanceUptimeSeconds/Minutes/Hours/Days  •  DatabaseName  •  DatabaseState  •  RecoveryModel\n"
               "DatabaseCreateDate  •  LastFullBackupTime  •  HoursSinceFullBackup  •  LastFullBackupStatus\n"
               "FullBackupAlertStatus  •  LastLogBackupTime  •  MinutesSinceLogBackup")
add_text_box(slide, Inches(0.8), Inches(5.9), Inches(12), Inches(1.5), schema_text, 12, False, MEDIUM_GRAY)

# ============================================================
# SLIDE 7: Workbook Dashboard
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "Workbook Dashboard — 4 Tabs", 36, True, WHITE)

tabs = [
    ("📊 Summary", [
        "4 KPI tiles: Instances, Databases,\nBackup Alerts, Compliance %",
        "Pie charts: Backup status distribution\nand Recovery model breakdown",
        "Instance uptime table with\nheat-map coloring"
    ], AZURE_BLUE),
    ("🖥️ Instances", [
        "Grid with online/offline status,\nSQL version, database count",
        "Uptime display in days/hours\nwith last-seen timestamp",
        "Line chart: uptime trend\nover time per instance"
    ], RGBColor(0x43, 0xA0, 0x47)),
    ("🗄️ Databases", [
        "Detailed grid with state,\nrecovery model, backup status",
        "Color-coded: ONLINE=green,\nFULL=blue, SIMPLE=orange",
        "Bar chart: database count\nper SQL instance"
    ], RGBColor(0xE6, 0x51, 0x00)),
    ("💾 Backups", [
        "Compliance table per instance\n(compliant, warning, critical)",
        "Alert grid: databases needing\nattention (sorted by severity)",
        "Log backup monitoring for\nFULL recovery model DBs"
    ], RGBColor(0x6A, 0x1B, 0x9A)),
]

for i, (tab_name, features, color) in enumerate(tabs):
    x = Inches(0.5 + i * 3.15)
    add_rounded_rect(slide, x, Inches(1.6), Inches(2.9), Inches(0.6), color, tab_name, 16, WHITE)
    for j, feature in enumerate(features):
        add_text_box(slide, x + Inches(0.15), Inches(2.4 + j * 1.4), Inches(2.6), Inches(1.3), feature, 12, False, DARK_GRAY)

# Parameters
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(12), Inches(0.3),
             "Interactive Parameters:  Subscription  •  Workspace  •  Time Range  •  SQL Instance (multi-select)  •  Database (multi-select)",
             14, True, DARK_BLUE)

# ============================================================
# SLIDE 8: Security
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "Security & Authentication", 36, True, WHITE)

# Two authentication options
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "Azure Authentication (Managed Identity)", 20, True, AZURE_BLUE)

azure_auth = [
    "• System-assigned MI on Automation Account",
    "• Arc VM MI for IMDS token (Hybrid Worker)",
    "• OAuth2 tokens for Azure Monitor & Key Vault",
    "• RBAC: Monitoring Metrics Publisher on DCR",
    "• RBAC: Key Vault Secrets User (SQL Auth only)",
    "• No credentials stored in runbook code"
]
add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3), azure_auth, 15, DARK_GRAY)

add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5),
             "SQL Server Authentication Options", 20, True, AZURE_BLUE)

# Windows Auth box
add_rounded_rect(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(1.8), LIGHT_BLUE)
add_text_box(slide, Inches(7.2), Inches(2.4), Inches(5.1), Inches(0.4), "Option A: Windows Authentication", 16, True, DARK_BLUE)
win_auth_desc = ("Best for domain-joined environments.\nHybrid Worker service account authenticates\nvia Kerberos/NTLM. No password management.\n"
                 "Requirements: Domain trust, SQL login for machine account.")
add_text_box(slide, Inches(7.2), Inches(2.8), Inches(5.1), Inches(1.2), win_auth_desc, 13, False, DARK_GRAY)

# SQL Auth box
add_rounded_rect(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(1.8), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(7.2), Inches(4.4), Inches(5.1), Inches(0.4), "Option B: SQL Authentication + Key Vault", 16, True, RGBColor(0xE6, 0x51, 0x00))
sql_auth_desc = ("Best for non-domain or mixed environments.\nCredentials stored securely in Azure Key Vault.\nManaged Identity retrieves secrets at runtime.\n"
                 "Requirements: Key Vault with SQL login secrets.")
add_text_box(slide, Inches(7.2), Inches(4.8), Inches(5.1), Inches(1.2), sql_auth_desc, 13, False, DARK_GRAY)

# SQL permissions
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.4),
             "Required SQL Server Permissions:", 16, True, DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.9), Inches(12), Inches(0.8),
             "VIEW SERVER STATE  •  VIEW ANY DATABASE  •  db_datareader on msdb (for backup history)",
             14, False, MEDIUM_GRAY)

# ============================================================
# SLIDE 9: Deployment Options
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), "Deployment — Easy as 1-2-3", 36, True, WHITE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(12), Inches(0.5),
             "Three ARM templates — deployable from the Azure Portal (no CLI required)", 20, False, AZURE_BLUE)

templates = [
    ("Step 1\nInfrastructure", "arm-template-infrastructure.json", [
        "Automation Account (System MI)",
        "Log Analytics Workspace",
        "Custom Table (20 columns)",
        "Key Vault (optional)"
    ], AZURE_BLUE),
    ("Step 2\nData Collection", "arm-template-data-collection.json", [
        "Data Collection Endpoint (DCE)",
        "Data Collection Rule (DCR)",
        "Stream declarations",
        "Transform KQL config"
    ], RGBColor(0x43, 0xA0, 0x47)),
    ("Step 3\nWorkbook", "arm-template-workbook.json", [
        "Azure Monitor Workbook",
        "4 interactive tabs",
        "Pre-configured KQL queries",
        "Color-coded visualizations"
    ], RGBColor(0x6A, 0x1B, 0x9A)),
]

for i, (title, filename, items, color) in enumerate(templates):
    x = Inches(0.8 + i * 4)
    add_rounded_rect(slide, x, Inches(2.3), Inches(3.5), Inches(0.9), color, title, 15, WHITE)
    add_text_box(slide, x + Inches(0.1), Inches(3.3), Inches(3.3), Inches(0.3), filename, 11, False, MEDIUM_GRAY, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.2), Inches(3.7 + j * 0.4), Inches(3.1), Inches(0.4),
                     f"✓  {item}", 13, False, DARK_GRAY)

# Manual steps
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.4),
             "Manual Steps (Portal-guided):", 16, True, DARK_BLUE)
manual = [
    "→  Configure RBAC (assign Monitoring Metrics Publisher role to Automation Account MI on DCR)",
    "→  Set up Hybrid Worker Group (add Arc-enabled server to Automation Account)",
    "→  Import & publish the runbook script (paste into Automation Account runbook editor)",
    "→  Create schedule and link to runbook with parameters (SQL instances, DCE endpoint, DCR ID)"
]
add_bullet_list(slide, Inches(0.8), Inches(5.9), Inches(12), Inches(1.5), manual, 13, DARK_GRAY)

# ============================================================
# SLIDE 10: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)
add_shape_bg(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.08), AZURE_BLUE)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             "Next Steps", 40, True, WHITE, PP_ALIGN.CENTER)

next_steps = [
    "1.  Review the Lab Deployment Guide (LabGuide-SQLServerMonitoring.md)",
    "2.  Deploy ARM templates to your subscription (Portal or script)",
    "3.  Configure Hybrid Worker on an Arc-enabled server or Azure VM",
    "4.  Import the runbook and create a monitoring schedule",
    "5.  Validate data in the workbook dashboard",
    "6.  Set up alert rules for backup SLA violations",
    "7.  Scale: add more SQL instances to the collection parameters"
]
add_bullet_list(slide, Inches(2), Inches(2.0), Inches(9), Inches(3.5), next_steps, 20, RGBColor(0xBB, 0xDE, 0xFB))

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
             "All materials provided:  Presentation  •  Architecture Diagram  •  Lab Guide  •  ARM Templates  •  Deployment Script  •  Runbook",
             16, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Thank You", 28, True, WHITE, PP_ALIGN.CENTER)

# Save to temp first (OneDrive paths can cause issues), then copy
temp_output = r"C:\temp\SQLServerMonitoring-Presentation.pptx"
final_output = os.path.join(
    r"C:\Users\rogeriob\OneDrive - Microsoft\repos\AzureMonitorAssets\Solutions\SQL Monitoring\Presentation",
    "SQLServerMonitoring-Presentation.pptx"
)
prs.save(temp_output)
import shutil
shutil.copy2(temp_output, final_output)
print(f"Presentation saved to: {final_output}")
print(f"Slides: {len(prs.slides)}")
